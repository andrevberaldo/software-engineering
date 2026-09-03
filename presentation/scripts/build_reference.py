"""Build reference.pptx: a themed template for pandoc's --reference-doc.

Palette (Midnight Executive + amber accent):
  navy   1E2761  - dominant, dark backgrounds / titles
  ice    CADCFC  - secondary, subtle text/backgrounds
  amber  E8A33D  - sharp accent, used sparingly
  ink    2B2B33  - body text on light backgrounds
  white  FFFFFF  - light backgrounds
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn, nsdecls
from lxml import etree

NAVY = "1E2761"
ICE = "CADCFC"
AMBER = "E8A33D"
INK = "2B2B33"
WHITE = "FFFFFF"

NSMAP_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

TITLE_TYPES = {
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
}


def set_bg(slide_or_layout, hexval):
    bg = slide_or_layout.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(hexval)


def set_placeholder_style(ph, hexval, size_pt=None, bold=None, name="Calibri", align=None):
    """Write formatting into the placeholder's list style (a:lstStyle/a:lvl1pPr/a:defRPr).

    This -- not a literal paragraph/run on the placeholder -- is what
    LibreOffice and PowerPoint actually use as the inherited default for
    slides built from this layout.
    """
    txBody = ph.text_frame._txBody
    old = txBody.find(qn("a:lstStyle"))
    if old is not None:
        txBody.remove(old)

    lvl_attrs = ' algn="%s"' % align if align else ""
    rpr_attrs = ""
    if size_pt is not None:
        rpr_attrs += ' sz="%d"' % int(size_pt.pt * 100)
    if bold is not None:
        rpr_attrs += ' b="%d"' % (1 if bold else 0)

    xml = (
        '<a:lstStyle %s>'
        '<a:lvl1pPr%s>'
        '<a:defRPr%s>'
        '<a:solidFill><a:srgbClr val="%s"/></a:solidFill>'
        '<a:latin typeface="%s"/>'
        '</a:defRPr>'
        '</a:lvl1pPr>'
        '</a:lstStyle>'
    ) % (nsdecls("a"), lvl_attrs, rpr_attrs, hexval, name)
    new_el = parse_xml(xml)
    bodyPr = txBody.find(qn("a:bodyPr"))
    bodyPr.addnext(new_el)


def edit_theme(prs):
    master = prs.slide_masters[0]
    theme_part = master.part.part_related_by(
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
    )
    root = etree.fromstring(theme_part.blob)
    clr_scheme = root.find(f".//{{{NSMAP_A}}}clrScheme")

    def set_srgb(tag, hexval):
        el = clr_scheme.find(f"{{{NSMAP_A}}}{tag}")
        for child in list(el):
            el.remove(child)
        srgb = etree.SubElement(el, "{%s}srgbClr" % NSMAP_A)
        srgb.set("val", hexval)

    set_srgb("dk1", INK)
    set_srgb("lt1", WHITE)
    set_srgb("dk2", NAVY)
    set_srgb("lt2", ICE)
    set_srgb("accent1", NAVY)
    set_srgb("accent2", ICE)
    set_srgb("accent3", AMBER)
    set_srgb("accent4", "3A4A9E")
    set_srgb("accent5", "8FA9E8")
    set_srgb("accent6", "6B7280")
    set_srgb("hlink", NAVY)
    set_srgb("folHlink", "6B7280")

    font_scheme = root.find(f".//{{{NSMAP_A}}}fontScheme")
    for major_minor in ("majorFont", "minorFont"):
        el = font_scheme.find(f"{{{NSMAP_A}}}{major_minor}")
        latin = el.find(f"{{{NSMAP_A}}}latin")
        latin.set("typeface", "Calibri")

    theme_part._blob = etree.tostring(
        root, xml_declaration=True, encoding="UTF-8", standalone=True
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    edit_theme(prs)

    layouts = {l.name: l for l in prs.slide_masters[0].slide_layouts}

    # --- Title Slide: dark navy, white title, ice subtitle ---
    lay = layouts["Title Slide"]
    set_bg(lay, NAVY)
    for ph in lay.placeholders:
        if ph.placeholder_format.type in TITLE_TYPES:
            set_placeholder_style(ph, WHITE, Pt(44), True)
        else:
            set_placeholder_style(ph, ICE, Pt(20), False)

    # --- Section Header: dark navy divider slides ---
    lay = layouts["Section Header"]
    set_bg(lay, NAVY)
    for ph in lay.placeholders:
        if ph.placeholder_format.type in TITLE_TYPES:
            set_placeholder_style(ph, WHITE, Pt(38), True)
        else:
            set_placeholder_style(ph, ICE, Pt(18), False)

    # --- Light content-style layouts ---
    for name in [
        "Title and Content",
        "Two Content",
        "Comparison",
        "Title Only",
        "Content with Caption",
        "Picture with Caption",
        "Blank",
    ]:
        lay = layouts[name]
        set_bg(lay, WHITE)
        for ph in lay.placeholders:
            ptype = ph.placeholder_format.type
            if ptype in TITLE_TYPES:
                set_placeholder_style(ph, NAVY, Pt(32), True)
            else:
                set_placeholder_style(ph, INK, Pt(16), False)

    set_bg(prs.slide_masters[0], WHITE)

    # python-pptx omits <p:sldIdLst> entirely when a deck has zero slides.
    # Pandoc's reference-doc handling expects that element to exist (even
    # empty) so it can inject its own slides -- without it, pandoc silently
    # drops every generated slide from presentation.xml. Add a throwaway
    # slide to force the element into existence, then remove just the
    # <p:sldId> entry (and its relationship/part) so zero slides remain.
    dummy = prs.slides.add_slide(layouts["Blank"])
    sld_id_lst = prs.slides._sldIdLst
    sld_id_el = list(sld_id_lst)[0]
    r_id = sld_id_el.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    sld_id_lst.remove(sld_id_el)
    prs.part.drop_rel(r_id)

    # Drop the printerSettings relationship too: pandoc's reference-doc
    # handling copies this relationship into generated decks but not the
    # binary part it points to, leaving a broken reference that fails
    # OOXML validation. Nothing in this template needs it.
    printer_rel_type = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/printerSettings"
    )
    for rid, rel in list(prs.part.rels.items()):
        if rel.reltype == printer_rel_type:
            prs.part.drop_rel(rid)

    out = "presentation/assets/reference.pptx"
    prs.save(out)
    print("saved", out)


if __name__ == "__main__":
    main()
